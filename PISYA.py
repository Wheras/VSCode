from pptx import Presentation
from pptx.util import Inches
import openpyxl
import os

# Функция очистки путей
def clean_path(path):
    return str(path).strip().strip('"') if path else ""

# Загружаем Excel
wb = openpyxl.load_workbook(r"D:/VSCode/AutoPowerPoint.xlsx")
ws = wb.active

# Создаём презентацию
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_layout = prs.slide_layouts[6]

# Проходим по строкам Excel
for i, row in enumerate(ws.iter_rows(min_row=2, values_only=True), start=2):
    if not any(row):
        print(f"⏭ Пропускаем пустую строку {i}")
        continue

    # Проверка длины строки (должно быть 12 значений)
    if len(row) < 12:
        print(f"⚠ Недостаточно данных в строке {i}, найдено {len(row)} — пропущено")
        continue

    # Распаковка и очистка
    (question_path, answer_path, background_path, slide_type,
     q_left, q_top, q_width, q_height,
     a_left, a_top, a_width, a_height) = row

    question_path = clean_path(question_path)
    answer_path = clean_path(answer_path)
    background_path = clean_path(background_path)
    slide_type = str(slide_type).strip()

    # Создаём слайд
    slide = prs.slides.add_slide(blank_layout)

    # Установка фона
    if background_path and os.path.isfile(background_path):
        slide.shapes.add_picture(background_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    else:
        print(f"⚠ Фон не найден в строке {i}: {background_path}")

    # Вставка вопроса
    if slide_type in ("Question", "Answer") and question_path and os.path.isfile(question_path):
        slide.shapes.add_picture(
            question_path,
            Inches(q_left), Inches(q_top),
            width=Inches(q_width), height=Inches(q_height)
        )
    else:
        print(f"⚠ Вопрос не вставлен (строка {i}): {question_path}")

    # Вставка ответа
    if slide_type == "Answer" and answer_path and os.path.isfile(answer_path):
        slide.shapes.add_picture(
            answer_path,
            Inches(a_left), Inches(a_top),
            width=Inches(a_width), height=Inches(a_height)
        )
    elif slide_type == "Answer":
        print(f"⚠ Ответ не вставлен (строка {i}): {answer_path}")

# Сохраняем
prs.save("presentation_with_question_answer.pptx")
print("✅ Презентация создана.")
